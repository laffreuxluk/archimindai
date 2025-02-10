from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import io
import os

def generate_pptx_report(df, predictions, anomalies, commentary, output_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ArchimindAI Report"
    slide.placeholders[1].text = "Analyse des données et résultats"

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Résumé des Données"
    content = f"Nombre d'enregistrements : {len(df)}\n"
    content += f"Nombre de variables : {df.shape[1]}\n"
    content += f"Nombre d'anomalies détectées : {len(anomalies)}"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    textbox.text_frame.text = content

    numeric_cols = df.select_dtypes(include=["float64", "int64"]).columns
    if len(numeric_cols) > 0:
        fig, ax = plt.subplots()
        ax.hist(df[numeric_cols[0]], bins=20)
        ax.set_title(f"Distribution de {numeric_cols[0]}")
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format="png")
        plt.close(fig)
        img_stream.seek(0)
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Histogramme de {numeric_cols[0]}"
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(6))
    
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Commentaires IA"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    textbox.text_frame.text = commentary if commentary else "Aucun commentaire généré."

    prs.save(output_path)
    print(f"Rapport PPTX sauvegardé dans {output_path}")

def generate_html_report(df, predictions, anomalies, commentary, output_path):
    html_content = "<html><head><meta charset='UTF-8'><title>ArchimindAI Report</title></head><body>"
    html_content += "<h1>ArchimindAI Report</h1>"
    html_content += f"<p>Nombre d'enregistrements : {len(df)}</p>"
    html_content += f"<p>Nombre de variables : {df.shape[1]}</p>"
    html_content += f"<p>Nombre d'anomalies détectées : {len(anomalies)}</p>"
    html_content += "<h2>Commentaires IA</h2>"
    html_content += f"<p>{commentary}</p>"
    html_content += "</body></html>"
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    print(f"Rapport HTML sauvegardé dans {output_path}")
